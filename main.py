"""
Parsedown — a fully local document-to-Markdown converter.

Everything runs on your machine. No external API calls, no cloud services,
no ML model downloads. The PDF path is the priority and is kept isolated from
the secondary converters (images / docx / pptx / xlsx) so it can never be
broken by them.

Conversions are remembered in a small local SQLite database (history.db) so
you can reopen, search, re-download, or delete past results — and it survives
restarting the app.

Run with:
    uvicorn main:app --reload
Then open http://127.0.0.1:8000
"""

from __future__ import annotations

import io
import sqlite3
import zipfile
from datetime import datetime
from pathlib import Path

from fastapi import FastAPI, UploadFile, File, Body
from fastapi.responses import JSONResponse, HTMLResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles

app = FastAPI(title="Parsedown", description="Local document → Markdown converter")

# Where the single-page frontend lives.
STATIC_DIR = Path(__file__).parent / "static"
# Local history database (created automatically on first run).
DB_PATH = Path(__file__).parent / "history.db"

# File extensions we know how to handle, grouped by converter.
PDF_EXTS = {".pdf"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp"}
DOCX_EXTS = {".docx"}
PPTX_EXTS = {".pptx"}
XLSX_EXTS = {".xlsx"}
SUPPORTED_EXTS = PDF_EXTS | IMAGE_EXTS | DOCX_EXTS | PPTX_EXTS | XLSX_EXTS


class ConversionError(Exception):
    """Raised by a converter when a file can't be processed.

    The message is safe to show directly to the user (friendly, one line).
    """


# --------------------------------------------------------------------------- #
# Markdown helpers
# --------------------------------------------------------------------------- #

def _md_escape_cell(value) -> str:
    """Make a cell value safe to drop into a Markdown table cell."""
    text = "" if value is None else str(value)
    return text.replace("\\", "\\\\").replace("|", "\\|").replace("\n", " ").strip()


def rows_to_md_table(rows: list[list]) -> str:
    """Render a list of rows (first row treated as header) as a Markdown table."""
    rows = [r for r in rows if any(c is not None and str(c).strip() != "" for c in r)]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    norm = [list(r) + [""] * (width - len(r)) for r in rows]
    header = norm[0]
    body = norm[1:]
    lines = ["| " + " | ".join(_md_escape_cell(c) for c in header) + " |"]
    lines.append("| " + " | ".join(["---"] * width) + " |")
    for r in body:
        lines.append("| " + " | ".join(_md_escape_cell(c) for c in r) + " |")
    return "\n".join(lines)


# --------------------------------------------------------------------------- #
# Converters — one function per file family. Each takes raw bytes and returns
# a Markdown string, or raises ConversionError with a user-friendly message.
# --------------------------------------------------------------------------- #

def convert_pdf(data: bytes) -> str:
    """Convert a PDF to Markdown using pymupdf4llm (top-priority path)."""
    import pymupdf  # PyMuPDF, a dependency of pymupdf4llm
    import pymupdf4llm

    try:
        # Open from memory so we never touch the disk.
        doc = pymupdf.open(stream=data, filetype="pdf")
    except Exception as exc:  # corrupt / not really a PDF
        raise ConversionError(
            "This file could not be opened as a PDF — it may be corrupt or "
            "password-protected."
        ) from exc

    try:
        markdown = pymupdf4llm.to_markdown(doc, show_progress=False)
    except Exception as exc:
        raise ConversionError(
            "Something went wrong while converting this PDF."
        ) from exc
    finally:
        doc.close()

    return markdown


def convert_image(data: bytes) -> str:
    """OCR an image into Markdown text using pytesseract.

    Degrades gracefully with an install hint if the Tesseract binary or the
    Python wrapper is missing.
    """
    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise ConversionError(
            "Image OCR needs the 'pytesseract' and 'Pillow' Python packages "
            "(pip install pytesseract Pillow)."
        ) from exc

    try:
        image = Image.open(io.BytesIO(data))
    except Exception as exc:
        raise ConversionError(
            "This image could not be opened — it may be corrupt."
        ) from exc

    try:
        text = pytesseract.image_to_string(image)
    except pytesseract.TesseractNotFoundError as exc:
        # The system Tesseract binary isn't installed — give the one-line hint.
        raise ConversionError(
            "Tesseract OCR isn't installed. Install it with "
            "'brew install tesseract' (macOS), 'apt install tesseract-ocr' "
            "(Linux), or the Windows installer, then try again."
        ) from exc
    except Exception as exc:
        raise ConversionError("OCR failed on this image.") from exc

    text = text.strip()
    if not text:
        return "_No text was detected in this image._"
    return text


def convert_docx(data: bytes) -> str:
    """Convert a Word .docx file to Markdown using mammoth (pure Python)."""
    try:
        import mammoth
    except ImportError as exc:
        raise ConversionError(
            "Word conversion needs the 'mammoth' package (pip install mammoth)."
        ) from exc

    try:
        result = mammoth.convert_to_markdown(io.BytesIO(data))
    except Exception as exc:
        raise ConversionError(
            "This .docx file could not be converted — it may be corrupt or "
            "not a real Word document."
        ) from exc

    return result.value


def convert_pptx(data: bytes) -> str:
    """Convert a PowerPoint .pptx file to Markdown using python-pptx.

    Each slide becomes a section; text frames become paragraphs and tables
    become Markdown tables.
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt  # noqa: F401  (imported to confirm install)
    except ImportError as exc:
        raise ConversionError(
            "PowerPoint conversion needs the 'python-pptx' package "
            "(pip install python-pptx)."
        ) from exc

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise ConversionError(
            "This .pptx file could not be opened — it may be corrupt or not a "
            "real PowerPoint file."
        ) from exc

    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {idx}")
        for shape in slide.shapes:
            # Tables
            if shape.has_table:
                rows = [
                    [cell.text for cell in row.cells] for row in shape.table.rows
                ]
                table_md = rows_to_md_table(rows)
                if table_md:
                    parts.append(table_md)
                continue
            # Text
            if shape.has_text_frame:
                text = "\n".join(
                    p.text for p in shape.text_frame.paragraphs if p.text.strip()
                ).strip()
                if text:
                    parts.append(text)
        parts.append("")  # blank line between slides

    markdown = "\n\n".join(p for p in parts if p != "").strip()
    return markdown or "_No text found in this presentation._"


def convert_xlsx(data: bytes) -> str:
    """Convert an Excel .xlsx workbook to Markdown — one table per sheet."""
    try:
        import openpyxl
    except ImportError as exc:
        raise ConversionError(
            "Excel conversion needs the 'openpyxl' package (pip install openpyxl)."
        ) from exc

    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        raise ConversionError(
            "This .xlsx file could not be opened — it may be corrupt or not a "
            "real Excel file."
        ) from exc

    parts: list[str] = []
    try:
        for ws in wb.worksheets:
            rows = [list(r) for r in ws.iter_rows(values_only=True)]
            table_md = rows_to_md_table(rows)
            parts.append(f"## {ws.title}")
            parts.append(table_md if table_md else "_Empty sheet._")
    finally:
        wb.close()

    return "\n\n".join(parts).strip() or "_This workbook has no data._"


def route_and_convert(filename: str, data: bytes) -> tuple[str, str]:
    """Auto-detect the file type by extension and dispatch to a converter.

    Returns (markdown, filetype_label).
    """
    ext = Path(filename).suffix.lower()
    if ext in PDF_EXTS:
        return convert_pdf(data), "PDF"
    if ext in IMAGE_EXTS:
        return convert_image(data), "Image"
    if ext in DOCX_EXTS:
        return convert_docx(data), "Word"
    if ext in PPTX_EXTS:
        return convert_pptx(data), "PowerPoint"
    if ext in XLSX_EXTS:
        return convert_xlsx(data), "Excel"
    raise ConversionError(
        f"Unsupported file type '{ext or 'unknown'}'. Supported: PDF, DOCX, "
        "PPTX, XLSX, and images (PNG/JPG/JPEG/TIFF/BMP/GIF/WEBP)."
    )


def build_stats(original_size: int, markdown: str) -> dict:
    """Compute the stats shown in the UI strip."""
    output_size = len(markdown.encode("utf-8"))
    reduction = 0.0
    if original_size > 0:
        reduction = round((1 - output_size / original_size) * 100, 1)
    # Rough token estimate: ~4 characters per token.
    est_tokens = len(markdown) // 4
    return {
        "original_size": original_size,
        "output_size": output_size,
        "reduction_pct": reduction,
        "est_tokens": est_tokens,
    }


# --------------------------------------------------------------------------- #
# Local history (SQLite)
# --------------------------------------------------------------------------- #

def get_db() -> sqlite3.Connection:
    """Open a connection to the local history database."""
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db() -> None:
    """Create the history table if it doesn't exist yet."""
    with get_db() as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS conversions (
                id            INTEGER PRIMARY KEY AUTOINCREMENT,
                filename      TEXT NOT NULL,
                filetype      TEXT NOT NULL,
                markdown      TEXT NOT NULL,
                original_size INTEGER NOT NULL,
                output_size   INTEGER NOT NULL,
                reduction_pct REAL NOT NULL,
                est_tokens    INTEGER NOT NULL,
                created_at    TEXT NOT NULL
            )
            """
        )


def save_conversion(filename: str, filetype: str, markdown: str, stats: dict) -> dict:
    """Persist one conversion and return the full record (including its id)."""
    created_at = datetime.now().isoformat(timespec="seconds")
    with get_db() as conn:
        cur = conn.execute(
            """
            INSERT INTO conversions
                (filename, filetype, markdown, original_size, output_size,
                 reduction_pct, est_tokens, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                filename,
                filetype,
                markdown,
                stats["original_size"],
                stats["output_size"],
                stats["reduction_pct"],
                stats["est_tokens"],
                created_at,
            ),
        )
        new_id = cur.lastrowid
    return {
        "id": new_id,
        "filename": filename,
        "filetype": filetype,
        "markdown": markdown,
        "stats": stats,
        "created_at": created_at,
    }


init_db()


# --------------------------------------------------------------------------- #
# Conversion API
# --------------------------------------------------------------------------- #

@app.post("/api/convert")
async def convert(file: UploadFile = File(...)):
    """Convert a single uploaded file to Markdown and save it to history.

    Always returns a clean JSON response — never crashes the server.
    """
    try:
        data = await file.read()
    except Exception:
        return JSONResponse(
            status_code=400,
            content={"error": "The uploaded file could not be read."},
        )

    if not data:
        return JSONResponse(
            status_code=400, content={"error": "The uploaded file is empty."}
        )

    try:
        markdown, filetype = route_and_convert(file.filename or "", data)
    except ConversionError as exc:
        return JSONResponse(status_code=400, content={"error": str(exc)})
    except Exception:
        # Last-resort safety net so a converter bug never takes down the server.
        return JSONResponse(
            status_code=500,
            content={"error": "An unexpected error occurred during conversion."},
        )

    stats = build_stats(len(data), markdown)
    record = save_conversion(file.filename or "file", filetype, markdown, stats)
    return record


@app.post("/api/convert-batch")
async def convert_batch(files: list[UploadFile] = File(...)):
    """Convert multiple files and return one result object per file.

    Each successful file is saved to history. Files that fail come back with an
    'error' field instead of markdown, so one bad file never sinks the batch.
    """
    results = []
    for file in files:
        name = file.filename or "file"
        try:
            data = await file.read()
            if not data:
                raise ConversionError("The file is empty.")
            markdown, filetype = route_and_convert(name, data)
            stats = build_stats(len(data), markdown)
            record = save_conversion(name, filetype, markdown, stats)
            results.append(record)
        except ConversionError as exc:
            results.append({"filename": name, "error": str(exc)})
        except Exception:
            results.append(
                {"filename": name, "error": "An unexpected error occurred."}
            )
    return {"results": results}


# --------------------------------------------------------------------------- #
# History API
# --------------------------------------------------------------------------- #

@app.get("/api/history")
async def list_history(q: str = ""):
    """List saved conversions (newest first), optionally filtered by a search.

    The search matches the filename or the converted Markdown text.
    """
    with get_db() as conn:
        if q.strip():
            like = f"%{q.strip()}%"
            rows = conn.execute(
                """
                SELECT id, filename, filetype, original_size, output_size,
                       reduction_pct, est_tokens, created_at
                FROM conversions
                WHERE filename LIKE ? OR markdown LIKE ?
                ORDER BY id DESC
                """,
                (like, like),
            ).fetchall()
        else:
            rows = conn.execute(
                """
                SELECT id, filename, filetype, original_size, output_size,
                       reduction_pct, est_tokens, created_at
                FROM conversions
                ORDER BY id DESC
                """
            ).fetchall()
    return {"items": [_row_to_summary(r) for r in rows]}


@app.get("/api/history/{item_id}")
async def get_history_item(item_id: int):
    """Return one full conversion (including its Markdown) by id."""
    with get_db() as conn:
        row = conn.execute(
            "SELECT * FROM conversions WHERE id = ?", (item_id,)
        ).fetchone()
    if row is None:
        return JSONResponse(status_code=404, content={"error": "Not found."})
    return _row_to_record(row)


@app.delete("/api/history/{item_id}")
async def delete_history_item(item_id: int):
    """Delete one saved conversion."""
    with get_db() as conn:
        conn.execute("DELETE FROM conversions WHERE id = ?", (item_id,))
    return {"ok": True}


@app.post("/api/history/clear")
async def clear_history():
    """Delete all saved conversions."""
    with get_db() as conn:
        conn.execute("DELETE FROM conversions")
    return {"ok": True}


@app.post("/api/export")
async def export_zip(ids: list[int] = Body(..., embed=True)):
    """Bundle the given saved conversions into a single .zip for download."""
    buffer = io.BytesIO()
    with get_db() as conn, zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        seen: dict[str, int] = {}
        for item_id in ids:
            row = conn.execute(
                "SELECT filename, markdown FROM conversions WHERE id = ?",
                (item_id,),
            ).fetchone()
            if row is None:
                continue
            stem = Path(row["filename"]).stem or "file"
            # Avoid name collisions inside the zip.
            count = seen.get(stem, 0)
            seen[stem] = count + 1
            name = f"{stem}.md" if count == 0 else f"{stem} ({count}).md"
            zf.writestr(name, row["markdown"])
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/zip",
        headers={"Content-Disposition": 'attachment; filename="parsedown.zip"'},
    )


def _row_to_summary(row: sqlite3.Row) -> dict:
    """Shape a DB row as a lightweight list item (no markdown body)."""
    return {
        "id": row["id"],
        "filename": row["filename"],
        "filetype": row["filetype"],
        "created_at": row["created_at"],
        "stats": {
            "original_size": row["original_size"],
            "output_size": row["output_size"],
            "reduction_pct": row["reduction_pct"],
            "est_tokens": row["est_tokens"],
        },
    }


def _row_to_record(row: sqlite3.Row) -> dict:
    """Shape a DB row as a full record (with markdown)."""
    rec = _row_to_summary(row)
    rec["markdown"] = row["markdown"]
    return rec


# --------------------------------------------------------------------------- #
# Frontend: serve the single-page app at "/" and its assets under /static.
# --------------------------------------------------------------------------- #

@app.get("/", response_class=HTMLResponse)
async def index():
    """Serve the single-page UI."""
    return (STATIC_DIR / "index.html").read_text(encoding="utf-8")


# Mounted last so it doesn't shadow the API/index routes above.
app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")
