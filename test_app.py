"""End-to-end smoke test for Parsedown.

Generates sample files (PDF, PPTX, XLSX), runs them through the live FastAPI app
via the in-process test client, and checks batch conversion, the SQLite history,
zip export, and the unsupported-file error path.

Run with:  python test_app.py
"""

import io
import sys

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from fastapi.testclient import TestClient

from main import app


def make_sample_pdf() -> bytes:
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    c.setFont("Helvetica-Bold", 18)
    c.drawString(72, 720, "Parsedown Sample Document")
    c.setFont("Helvetica", 12)
    c.drawString(72, 690, "This is a sample paragraph used to verify conversion.")
    c.save()
    return buf.getvalue()


def make_sample_pptx() -> bytes:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Update"
    slide.placeholders[1].text = "Revenue is up and morale is high."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_sample_xlsx() -> bytes:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Region", "Units"])
    ws.append(["North", 120])
    ws.append(["South", 95])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def main() -> int:
    client = TestClient(app)
    failures = []

    def check(label, cond, detail=""):
        if cond:
            print(f"[OK]   {label}")
        else:
            print(f"[FAIL] {label} {detail}")
            failures.append(f"{label} {detail}")

    # 1. Server boots and serves the UI.
    r = client.get("/")
    check("Server boots and serves UI", r.status_code == 200 and "Parsedown" in r.text)

    # 2. PDF converts and is saved with an id + stats.
    r = client.post("/api/convert",
                    files={"file": ("sample.pdf", make_sample_pdf(), "application/pdf")})
    pdf_ok = r.status_code == 200 and "Parsedown Sample Document" in r.json().get("markdown", "")
    check("PDF converts to Markdown", pdf_ok, "" if pdf_ok else r.text)
    if pdf_ok:
        s = r.json()["stats"]
        print(f"        stats: orig={s['original_size']}B out={s['output_size']}B "
              f"reduction={s['reduction_pct']}% ~{s['est_tokens']} tokens, id={r.json()['id']}")

    # 3. PPTX converts.
    r = client.post("/api/convert",
                    files={"file": ("deck.pptx", make_sample_pptx(),
                                    "application/vnd.openxmlformats-officedocument.presentationml.presentation")})
    check("PPTX converts to Markdown",
          r.status_code == 200 and "Quarterly Update" in r.json().get("markdown", ""),
          "" if r.status_code == 200 else r.text)

    # 4. XLSX converts to a Markdown table.
    r = client.post("/api/convert",
                    files={"file": ("book.xlsx", make_sample_xlsx(),
                                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")})
    xlsx_md = r.json().get("markdown", "") if r.status_code == 200 else ""
    check("XLSX converts to a Markdown table",
          "| Region | Units |" in xlsx_md and "North" in xlsx_md,
          "" if r.status_code == 200 else r.text)

    # 5. Batch conversion returns one result per file (incl. a failing one).
    r = client.post("/api/convert-batch", files=[
        ("files", ("a.pdf", make_sample_pdf(), "application/pdf")),
        ("files", ("b.xlsx", make_sample_xlsx(), "application/octet-stream")),
        ("files", ("bad.xyz", b"nonsense", "application/octet-stream")),
    ])
    batch = r.json().get("results", []) if r.status_code == 200 else []
    ok_count = sum(1 for x in batch if not x.get("error"))
    fail_count = sum(1 for x in batch if x.get("error"))
    check("Batch returns per-file results", len(batch) == 3 and ok_count == 2 and fail_count == 1,
          f"(ok={ok_count}, failed={fail_count})")
    batch_ids = [x["id"] for x in batch if not x.get("error")]

    # 6. History lists saved conversions and search works.
    r = client.get("/api/history")
    items = r.json().get("items", [])
    check("History lists saved conversions", r.status_code == 200 and len(items) >= 4)

    r = client.get("/api/history", params={"q": "Quarterly"})
    check("History search finds text inside Markdown",
          r.status_code == 200 and any("deck" in i["filename"] for i in r.json().get("items", [])))

    # 7. Export selected ids as a zip.
    if batch_ids:
        r = client.post("/api/export", json={"ids": batch_ids})
        is_zip = r.status_code == 200 and r.content[:2] == b"PK"
        check("Export bundles selected items into a .zip", is_zip)

    # 8. Delete one history item.
    if items:
        del_id = items[0]["id"]
        r = client.delete(f"/api/history/{del_id}")
        gone = client.get(f"/api/history/{del_id}").status_code == 404
        check("Delete removes a history item", r.status_code == 200 and gone)

    # 9. Unsupported file -> clean 400, server stays up.
    r = client.post("/api/convert",
                    files={"file": ("notes.xyz", b"random", "application/octet-stream")})
    check("Unsupported file returns a clean error",
          r.status_code == 400 and "Unsupported" in r.json().get("error", ""))

    print("-" * 56)
    if failures:
        print(f"RESULT: {len(failures)} FAILED")
        return 1
    print("RESULT: ALL TESTS PASSED")
    return 0


if __name__ == "__main__":
    sys.exit(main())
